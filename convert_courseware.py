"""把 PDF / PPT 课件转成 courseware 用的 Markdown（front-matter + 分节 + 术语表）。

用法：
    python convert_courseware.py 课件.pdf -o courseware/课程名.md [--course 课程名]
    python convert_courseware.py 课件.pptx --no-glossary

`--glossary`（默认开）会用 DeepSeek 把提取到的英文术语配上中文，写成 `## 术语表`。
"""

from __future__ import annotations

import argparse
import os
import re
import sys


def extract_pdf(path: str):
    import fitz  # pymupdf
    doc = fitz.open(path)
    pages = []
    for i, page in enumerate(doc, 1):
        text = page.get_text("text").strip()
        if text:
            pages.append((f"第 {i} 页", text))
    doc.close()
    return pages


def extract_pptx(path: str):
    from pptx import Presentation
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(r.text for r in para.runs).strip()
                    if t:
                        parts.append(t)
            if getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        parts.append(" | ".join(cells))
        text = "\n".join(parts).strip()
        if text:
            slides.append((f"第 {i} 页", text))
    return slides


def _md_pages(pages) -> str:
    out = []
    for title, body in pages:
        out.append(f"## {title}\n")
        out.append(body.rstrip())
        out.append("")
    return "\n".join(out)


def _fill_glossary(md_text: str, course_name: str) -> str:
    """用 DeepSeek 给文本里的术语配中文，生成 `## 术语表` 表格。"""
    try:
        import courseware, translation
        # 先写成一个临时结构再解析术语（避免依赖文件路径）
        import tempfile
        with tempfile.NamedTemporaryFile("w", suffix=".md", delete=False, encoding="utf-8") as f:
            f.write(md_text)
            tmp = f.name
        course = courseware.load_course(tmp)
        os.unlink(tmp)
        terms = course.terms
        if not terms:
            return md_text
        tr = translation.Translator()
        prompt = ("把以下英文术语翻译成简体中文，每行输出 `英文 = 中文`，不要其它内容：" + "\n".join(terms))
        raw = tr.translate(prompt)
        pairs = {}
        for line in raw.splitlines():
            line = line.strip()
            if "=" in line:
                en, zh = line.split("=", 1)
                en, zh = en.strip(), zh.strip()
                if en and zh and en.lower() in {t.lower() for t in terms}:
                    pairs[en] = zh
        # 用解析到的原文术语（保序）
        table = ["| English | 中文 |", "| --- | --- |"]
        for t in terms:
            zh = pairs.get(t) or pairs.get(next((k for k in pairs if k.lower() == t.lower()), ""), "")
            table.append(f"| {t} | {zh} |")
        gloss = "## 术语表\n\n" + "\n".join(table) + "\n\n"
        # 插到 front-matter 之后
        if md_text.startswith("---"):
            idx = md_text.find("\n---\n")
            if idx >= 0:
                return md_text[:idx + 5] + "\n\n" + gloss + md_text[idx + 5:].lstrip("\n")
        return gloss + md_text
    except Exception as e:
        try:
            print(f"[警告] 术语表生成失败：{type(e).__name__}: {e}", file=sys.stderr)
        except Exception:
            pass
        return md_text


def main():
    ap = argparse.ArgumentParser(description="PDF/PPT 课件 -> 可用的 courseware Markdown")
    ap.add_argument("input", help="pdf / pptx 路径")
    ap.add_argument("-o", "--output", default=None, help="输出 md 路径（默认 courseware/<输入名>.md）")
    ap.add_argument("--course", default=None, help="课程名（默认取文件名）")
    ap.add_argument("--no-glossary", action="store_true", help="不自动生成术语表")
    args = ap.parse_args()

    path = args.input
    ext = os.path.splitext(path)[1].lower()
    if ext in (".pdf",):
        pages = extract_pdf(path)
    elif ext in (".pptx", ".ppt"):
        pages = extract_pptx(path)
    else:
        print("只支持 .pdf / .pptx / .ppt"); return

    if not pages:
        print("未提取到文本内容。"); return

    name = args.course or os.path.splitext(os.path.basename(path))[0]
    md = f"---\ncourse: {name}\n---\n\n" + _md_pages(pages)
    if not args.no_glossary:
        md = _fill_glossary(md, name)

    out = args.output or os.path.join(os.path.dirname(os.path.abspath(__file__)),
                                      "courseware", f"{name}.md")
    os.makedirs(os.path.dirname(out), exist_ok=True)
    with open(out, "w", encoding="utf-8") as f:
        f.write(md)
    print("已生成：", out)
    print("字符数：", len(md), "| 页数：", len(pages))


if __name__ == "__main__":
    main()
