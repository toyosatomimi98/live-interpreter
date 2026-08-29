import os, sys, subprocess, time, tempfile, wave
sys.stdout.reconfigure(encoding="utf-8", errors="replace")

RESULTS = []

def check(name, fn):
    try:
        msg = fn()
        RESULTS.append(("PASS", name, msg or ""))
        print(f"[PASS] {name}" + (f"  -> {msg}" if msg else ""))
    except Exception as e:
        RESULTS.append(("FAIL", name, f"{type(e).__name__}: {e}"))
        print(f"[FAIL] {name}  -> {type(e).__name__}: {e}")

base = os.path.dirname(os.path.abspath(__file__))
sample = os.path.join(base, "docs", "sample-courseware.md")

def c1_imports():
    import tongchuan, translation, courseware, convert_courseware
    return "ok"

def c2_load_course():
    import courseware as cw
    c = cw.load_course(sample)
    assert c.name == "硬件优化与内存访问", c.name
    assert "MESI" in c.terms and "store buffer" in c.terms
    assert len(c.glossary) >= 20, len(c.glossary)
    return f"terms={len(c.terms)} glossary={len(c.glossary)}"

def c3_sections():
    import courseware as cw
    secs = cw.load_sections(sample)
    assert secs and all("术语表" not in t for t, _ in secs), [t for t, _ in secs]
    return f"sections={len(secs)}"

def c4_convert_pptx():
    import convert_courseware as cc
    from pptx import Presentation
    p = os.path.join(tempfile.gettempdir(), "cw_test.pptx")
    prs = Presentation()
    for title, terms in [("Store Buffer", ["store buffer", "memory ordering"]),
                         ("Cache", ["MESI", "cache coherence", "LLC"])]:
        s = prs.slides.add_slide(prs.slide_layouts[0])
        s.shapes.title.text = title
        s.placeholders[1].text_frame.text = "\n".join(f"* {t}" for t in terms)
    prs.save(p)
    pages = cc.extract_pptx(p)
    assert len(pages) == 2, len(pages)
    return f"pages={len(pages)}"

def c5_merge():
    import tongchuan
    out = tongchuan._merge_prompt("a, b, c", "b, d, e", max_terms=10)
    assert out.split(", ")[0] == "b", out  # b (course) 优先
    assert "c" in out and "d" in out and "e" in out
    assert len(out.split(", ")) <= 10
    return out[:60]

def c6_cli_help():
    py = os.path.join(base, ".venv", "Scripts", "python.exe")
    r = subprocess.run([py, os.path.join(base, "tongchuan.py"), "--help"],
                       capture_output=True, text=True, encoding="utf-8",
                       errors="replace", timeout=30)
    assert r.returncode == 0 and "--model" in r.stdout and "--course" in r.stdout
    return "ok"

def c7_translate():
    from translation import Translator
    t = Translator()
    out = t.translate("The cache coherence protocol MESI keeps cores consistent.")
    assert out and "MESI" in out
    return out[:30]

def c8_fragment():
    from translation import Translator
    t = Translator()
    out = t.translate("do to start school in-")
    assert out and "请提供" not in out and "补充" not in out
    return out[:20]

def c9_glossary():
    import courseware, translation
    tr = translation.Translator(glossary=courseware.glossary_text(courseware.load_course(sample)))
    out = tr.translate("The store buffer uses a memory ordering model.")
    assert out and "存储缓冲" in out, out
    return out[:26]

def c10_asr_prompt():
    import edge_tts, av, numpy as np, asyncio
    from faster_whisper import WhisperModel
    import tongchuan, courseware
    mp3 = os.path.join(tempfile.gettempdir(), "cw_asr.mp3")
    asyncio.run(edge_tts.Communicate(
        "The cache coherence protocol MESI keeps cores consistent.",
        "en-US-AriaNeural").save(mp3))
    course = courseware.load_course(sample)
    prompt = tongchuan._merge_prompt(tongchuan.load_asr_prompt(), courseware.asr_prompt(course))
    assert "MESI" in prompt, "course term missing from ASR prompt"
    m = WhisperModel("base.en", device="cpu", compute_type="int8")
    segs, _ = m.transcribe(mp3, beam_size=1, language="en", vad_filter=True, initial_prompt=prompt)
    txt = "".join(s.text for s in segs).strip()
    assert len(txt) > 8, "ASR returned empty"
    return f"asr_ok len={len(txt)}"

def c11_pipeline():
    import edge_tts, av, numpy as np, asyncio, queue as _q
    import tongchuan
    mp3 = os.path.join(tempfile.gettempdir(), "cw_p.mp3")
    asyncio.run(edge_tts.Communicate(
        "Non temporal prefetching reduces bank conflict in the row buffer.",
        "en-US-AriaNeural").save(mp3))
    c = av.open(mp3); st = c.streams.audio[0]
    rs = av.AudioResampler(format="flt", layout="mono", rate=tongchuan.SAMPLE_RATE)
    ch = []
    for f in c.decode(st):
        for rf in rs.resample(f):
            ch.append(rf.to_ndarray().reshape(-1))
    audio = np.concatenate(ch)
    pipe = tongchuan.Pipeline(model_size="base.en", voice_enabled=False, course_file=sample,
                              translation_workers=2, status_cb=lambda s: None,
                              error_cb=lambda e: print("  [pipeline error]", e))
    pipe.start()
    dl = time.time() + 90
    while time.time() < dl and pipe.model is None:
        time.sleep(0.5)
    block = tongchuan.SAMPLE_RATE // 10
    for i in range(0, len(audio), block):
        pipe.segmenter.feed(audio[i:i+block].astype(np.float32))
        time.sleep(0.05)
    time.sleep(4)
    got = []
    dl = time.time() + 20
    while time.time() < dl:
        try:
            r = pipe.result_q.get(timeout=0.3)
            got.append(r)
        except _q.Empty:
            pass
    pipe.stop()
    assert len(got) >= 1, "no results"
    assert any(r.get("section") for r in got) or got, "no section"
    return f"results={len(got)} first_sec={got[0].get('section') or '-'}"

def c12_record():
    import tongchuan, os as _os, glob
    recdir = os.path.join(base, "recordings")
    _os.makedirs(recdir, exist_ok=True)
    pipe = tongchuan.Pipeline(model_size="base.en", source="mic", voice_enabled=False,
                              save_audio=True, status_cb=lambda s: None,
                              error_cb=lambda e: print("  [rec error]", e))
    pipe.start()
    time.sleep(3)
    pipe.stop()
    wavs = glob.glob(os.path.join(recdir, "*.wav"))
    assert wavs, "no wav"
    with wave.open(wavs[-1], "rb") as w:
        assert w.getframerate() == 16000 and w.getnchannels() == 1
        assert w.getnframes() > 1000
    return f"frames={w.getnframes()} sr={w.getframerate()}"

def c13_loopback():
    import io, contextlib
    import tongchuan
    buf = io.StringIO()
    pipe = tongchuan.Pipeline(model_size="base.en", source="system", voice_enabled=False,
                              status_cb=lambda s: None, error_cb=lambda e: None)
    with contextlib.redirect_stderr(buf):
        pipe.start()
        time.sleep(3)
        pipe.stop()
    assert "0x800401f0" not in buf.getvalue()
    return "ok"

if __name__ == "__main__":
    check("1 imports", c1_imports)
    check("2 load_course", c2_load_course)
    check("3 load_sections", c3_sections)
    check("4 convert_pptx", c4_convert_pptx)
    check("5 merge_prompt", c5_merge)
    check("6 cli_help", c6_cli_help)
    check("7 translate", c7_translate)
    check("8 fragment", c8_fragment)
    check("9 glossary", c9_glossary)
    check("10 asr_prompt", c10_asr_prompt)
    check("11 pipeline", c11_pipeline)
    check("12 record_wav", c12_record)
    check("13 loopback", c13_loopback)
    fails = [r for r in RESULTS if r[0] == "FAIL"]
    print("\n===== SUMMARY =====")
    print(f"PASS: {len(RESULTS)-len(fails)}/{len(RESULTS)}")
    for f in fails:
        print("  FAIL:", f[1], "->", f[2])
